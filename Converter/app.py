"""
FileForge — Production-Ready File Conversion API
Flask backend supporting 7+ conversion types with security, validation, and logging.
"""

import os
import sys
import uuid
import shutil
import logging
import tempfile
import subprocess
from pathlib import Path
from datetime import datetime

from flask import Flask, request, send_file, jsonify, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader, PdfWriter, PdfMerger

# ─── App Setup ────────────────────────────────────────────────────────────────

BASE_DIR = Path(__file__).resolve().parent
UPLOAD_DIR = BASE_DIR / 'uploads'
CONVERTED_DIR = BASE_DIR / 'converted'
UPLOAD_DIR.mkdir(exist_ok=True)
CONVERTED_DIR.mkdir(exist_ok=True)

app = Flask(__name__, static_folder='static', static_url_path='')
app.config['MAX_CONTENT_LENGTH'] = 20 * 1024 * 1024  # 20 MB
CORS(app)

# ─── Logging ──────────────────────────────────────────────────────────────────

LOG_DIR = BASE_DIR / 'logs'
LOG_DIR.mkdir(exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[
        logging.FileHandler(LOG_DIR / 'fileforge.log'),
        logging.StreamHandler(sys.stdout),
    ],
)
logger = logging.getLogger('FileForge')

# ─── Helpers ──────────────────────────────────────────────────────────────────

ALLOWED_MIME = {
    'pdf':   ['application/pdf'],
    'docx':  ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword'],
    'pptx':  ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint'],
    'image': ['image/jpeg', 'image/png', 'image/webp', 'image/bmp', 'image/tiff'],
}


def _unique_name(ext: str) -> str:
    """Generate a unique filename with the given extension."""
    return f"{uuid.uuid4().hex}{ext}"


def _save_upload(file_storage, expected_type: str = None) -> Path:
    """Save an uploaded file to a unique temp path and return its path."""
    original = secure_filename(file_storage.filename) or 'upload'
    ext = Path(original).suffix.lower()
    dest = UPLOAD_DIR / _unique_name(ext)
    file_storage.save(str(dest))
    logger.info(f"Saved upload: {original} → {dest.name} ({dest.stat().st_size} bytes)")
    return dest


def _cleanup(*paths):
    """Silently remove files."""
    for p in paths:
        try:
            if p and Path(p).exists():
                Path(p).unlink()
        except Exception:
            pass


def _validate_extension(filename: str, allowed_exts: list) -> bool:
    ext = Path(filename).suffix.lower()
    return ext in allowed_exts


def _is_scanned_pdf(pdf_path: str) -> bool:
    """Heuristic check: if a PDF has no extractable text, it's likely scanned."""
    try:
        reader = PdfReader(str(pdf_path))
        total_text = ''
        for page in reader.pages[:5]:  # Check first 5 pages
            text = page.extract_text() or ''
            total_text += text.strip()
        return len(total_text) < 10  # Less than 10 chars = probably scanned
    except Exception:
        return False


def _find_libreoffice():
    """Find LibreOffice executable on the system."""
    if sys.platform == 'win32':
        candidates = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        ]
        for c in candidates:
            if os.path.isfile(c):
                return c
    else:
        for cmd in ['soffice', 'libreoffice']:
            if shutil.which(cmd):
                return cmd
    return None


LIBREOFFICE = _find_libreoffice()

# ─── Routes: Static ──────────────────────────────────────────────────────────

@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')


@app.route('/health')
def health():
    return jsonify({
        'status': 'ok',
        'timestamp': datetime.utcnow().isoformat(),
        'libreoffice': LIBREOFFICE is not None,
    })


# ─── Route: PDF → DOCX ──────────────────────────────────────────────────────

@app.route('/api/convert/pdf-to-docx', methods=['POST'])
def pdf_to_docx():
    input_path = output_path = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not _validate_extension(file.filename, ['.pdf']):
            return jsonify({'error': 'Invalid file type. Expected: .pdf'}), 400

        input_path = _save_upload(file, 'pdf')

        # Check for scanned PDF
        if _is_scanned_pdf(input_path):
            _cleanup(input_path)
            return jsonify({'error': 'This appears to be a scanned PDF. OCR is required but not supported yet.'}), 422

        output_path = CONVERTED_DIR / _unique_name('.docx')

        from pdf2docx import Converter
        cv = Converter(str(input_path))
        cv.convert(str(output_path))
        cv.close()

        if not output_path.exists() or output_path.stat().st_size == 0:
            return jsonify({'error': 'Conversion produced an empty file'}), 500

        logger.info(f"PDF→DOCX conversion successful: {output_path.name}")
        original_stem = Path(secure_filename(file.filename)).stem
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name=f"{original_stem}.docx",
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        )
    except ImportError:
        return jsonify({'error': 'pdf2docx library not installed on server'}), 500
    except Exception as e:
        logger.error(f"PDF→DOCX error: {e}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        _cleanup(input_path, output_path)


# ─── Route: DOCX → PDF ──────────────────────────────────────────────────────

@app.route('/api/convert/docx-to-pdf', methods=['POST'])
def docx_to_pdf():
    input_path = output_path = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not _validate_extension(file.filename, ['.docx', '.doc']):
            return jsonify({'error': 'Invalid file type. Expected: .docx or .doc'}), 400

        input_path = _save_upload(file, 'docx')

        if sys.platform == 'win32':
            # Use docx2pdf on Windows (uses MS Word COM)
            try:
                from docx2pdf import convert as docx2pdf_convert
                output_path = CONVERTED_DIR / _unique_name('.pdf')
                docx2pdf_convert(str(input_path), str(output_path))
            except ImportError:
                if LIBREOFFICE:
                    output_path = _convert_with_libreoffice(input_path, 'pdf')
                else:
                    return jsonify({'error': 'Neither docx2pdf nor LibreOffice available'}), 500
        else:
            if not LIBREOFFICE:
                return jsonify({'error': 'LibreOffice not installed on server'}), 500
            output_path = _convert_with_libreoffice(input_path, 'pdf')

        if not output_path or not output_path.exists() or output_path.stat().st_size == 0:
            return jsonify({'error': 'Conversion produced an empty file'}), 500

        logger.info(f"DOCX→PDF conversion successful: {output_path.name}")
        original_stem = Path(secure_filename(file.filename)).stem
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name=f"{original_stem}.pdf",
            mimetype='application/pdf',
        )
    except Exception as e:
        logger.error(f"DOCX→PDF error: {e}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        _cleanup(input_path, output_path)


# ─── Route: PPT → PDF ────────────────────────────────────────────────────────

@app.route('/api/convert/ppt-to-pdf', methods=['POST'])
def ppt_to_pdf():
    input_path = output_path = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not _validate_extension(file.filename, ['.pptx', '.ppt']):
            return jsonify({'error': 'Invalid file type. Expected: .pptx or .ppt'}), 400

        input_path = _save_upload(file, 'pptx')

        if not LIBREOFFICE:
            return jsonify({'error': 'LibreOffice not installed. PPT→PDF requires LibreOffice.'}), 500

        output_path = _convert_with_libreoffice(input_path, 'pdf')

        if not output_path or not output_path.exists() or output_path.stat().st_size == 0:
            return jsonify({'error': 'Conversion produced an empty file'}), 500

        logger.info(f"PPT→PDF conversion successful: {output_path.name}")
        original_stem = Path(secure_filename(file.filename)).stem
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name=f"{original_stem}.pdf",
            mimetype='application/pdf',
        )
    except Exception as e:
        logger.error(f"PPT→PDF error: {e}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        _cleanup(input_path, output_path)


# ─── Route: PDF → Image ──────────────────────────────────────────────────────

@app.route('/api/convert/pdf-to-image', methods=['POST'])
def pdf_to_image():
    input_path = None
    output_paths = []
    zip_path = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not _validate_extension(file.filename, ['.pdf']):
            return jsonify({'error': 'Invalid file type. Expected: .pdf'}), 400

        input_path = _save_upload(file, 'pdf')

        try:
            from pdf2image import convert_from_path
        except ImportError:
            return jsonify({'error': 'pdf2image library not installed on server'}), 500

        # Find poppler path on Windows
        poppler_path = None
        if sys.platform == 'win32':
            # Check common locations
            for p in [r'C:\poppler\Library\bin', r'C:\poppler\bin',
                      r'C:\Program Files\poppler\Library\bin',
                      r'C:\Program Files\poppler\bin']:
                if os.path.isdir(p):
                    poppler_path = p
                    break

        kwargs = {}
        if poppler_path:
            kwargs['poppler_path'] = poppler_path

        try:
            images = convert_from_path(str(input_path), dpi=200, **kwargs)
        except Exception as exc:
            if 'poppler' in str(exc).lower() or 'pdftoppm' in str(exc).lower():
                return jsonify({
                    'error': 'Poppler is not installed. PDF→Image requires Poppler. See README for setup instructions.'
                }), 500
            raise

        if len(images) == 1:
            output_path = CONVERTED_DIR / _unique_name('.jpg')
            images[0].save(str(output_path), 'JPEG', quality=95)
            original_stem = Path(secure_filename(file.filename)).stem
            logger.info(f"PDF→JPG conversion successful: {output_path.name}")
            resp = send_file(
                str(output_path),
                as_attachment=True,
                download_name=f"{original_stem}.jpg",
                mimetype='image/jpeg',
            )
            output_paths.append(output_path)
            return resp
        else:
            # Multiple pages → ZIP
            import zipfile
            zip_path = CONVERTED_DIR / _unique_name('.zip')
            original_stem = Path(secure_filename(file.filename)).stem
            with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
                for i, img in enumerate(images):
                    img_path = CONVERTED_DIR / _unique_name('.jpg')
                    img.save(str(img_path), 'JPEG', quality=95)
                    zf.write(str(img_path), f"{original_stem}_page_{i+1}.jpg")
                    output_paths.append(img_path)

            logger.info(f"PDF→JPG (multi-page) conversion successful: {zip_path.name}")
            return send_file(
                str(zip_path),
                as_attachment=True,
                download_name=f"{original_stem}_images.zip",
                mimetype='application/zip',
            )
    except Exception as e:
        logger.error(f"PDF→Image error: {e}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        _cleanup(input_path, zip_path, *output_paths)


# ─── Route: Image → PDF ──────────────────────────────────────────────────────

@app.route('/api/convert/image-to-pdf', methods=['POST'])
def image_to_pdf():
    input_paths = []
    output_path = None
    try:
        files = request.files.getlist('files')
        if not files:
            return jsonify({'error': 'No files uploaded'}), 400

        for f in files:
            if not _validate_extension(f.filename, ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff']):
                return jsonify({'error': f'Invalid file type: {f.filename}. Expected image file.'}), 400
            input_paths.append(_save_upload(f, 'image'))

        output_path = CONVERTED_DIR / _unique_name('.pdf')

        from PIL import Image
        images = []
        for p in input_paths:
            img = Image.open(str(p))
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            images.append(img)

        if images:
            if len(images) == 1:
                images[0].save(str(output_path), 'PDF', resolution=150)
            else:
                images[0].save(str(output_path), 'PDF', resolution=150, save_all=True, append_images=images[1:])

        if not output_path.exists() or output_path.stat().st_size == 0:
            return jsonify({'error': 'Conversion produced an empty file'}), 500

        logger.info(f"Image→PDF conversion successful: {len(images)} images → {output_path.name}")
        original_stem = Path(secure_filename(files[0].filename)).stem
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name=f"{original_stem}.pdf",
            mimetype='application/pdf',
        )
    except Exception as e:
        logger.error(f"Image→PDF error: {e}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        _cleanup(*input_paths)
        _cleanup(output_path)


# ─── Route: PDF → PPT ────────────────────────────────────────────────────────

@app.route('/api/convert/pdf-to-ppt', methods=['POST'])
def pdf_to_ppt():
    input_path = output_path = None
    img_paths = []
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not _validate_extension(file.filename, ['.pdf']):
            return jsonify({'error': 'Invalid file type. Expected: .pdf'}), 400

        input_path = _save_upload(file, 'pdf')
        output_path = CONVERTED_DIR / _unique_name('.pptx')

        # Convert PDF pages to images, then embed in PPTX slides
        try:
            from pdf2image import convert_from_path
        except ImportError:
            return jsonify({'error': 'pdf2image library not installed on server'}), 500

        poppler_path = None
        if sys.platform == 'win32':
            for p in [r'C:\poppler\Library\bin', r'C:\poppler\bin',
                      r'C:\Program Files\poppler\Library\bin']:
                if os.path.isdir(p):
                    poppler_path = p
                    break

        kwargs = {}
        if poppler_path:
            kwargs['poppler_path'] = poppler_path

        try:
            images = convert_from_path(str(input_path), dpi=200, **kwargs)
        except Exception as exc:
            if 'poppler' in str(exc).lower():
                return jsonify({'error': 'Poppler is not installed. PDF→PPT requires Poppler.'}), 500
            raise

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            img_path = CONVERTED_DIR / _unique_name('.jpg')
            img.save(str(img_path), 'JPEG', quality=95)
            img_paths.append(img_path)
            slide.shapes.add_picture(str(img_path), Inches(0), Inches(0),
                                     prs.slide_width, prs.slide_height)

        prs.save(str(output_path))

        if not output_path.exists() or output_path.stat().st_size == 0:
            return jsonify({'error': 'Conversion produced an empty file'}), 500

        logger.info(f"PDF→PPT conversion successful: {output_path.name}")
        original_stem = Path(secure_filename(file.filename)).stem
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name=f"{original_stem}.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        )
    except ImportError as ie:
        return jsonify({'error': f'Required library not installed: {ie}'}), 500
    except Exception as e:
        logger.error(f"PDF→PPT error: {e}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        _cleanup(input_path, output_path, *img_paths)


# ─── Route: Merge PDF ────────────────────────────────────────────────────────

@app.route('/api/convert/merge-pdf', methods=['POST'])
def merge_pdf():
    input_paths = []
    output_path = None
    try:
        files = request.files.getlist('files')
        if not files or len(files) < 2:
            return jsonify({'error': 'Please upload at least 2 PDF files to merge'}), 400

        for f in files:
            if not _validate_extension(f.filename, ['.pdf']):
                return jsonify({'error': f'File "{f.filename}" is not a PDF'}), 400
            input_paths.append(_save_upload(f, 'pdf'))

        output_path = CONVERTED_DIR / _unique_name('.pdf')
        merger = PdfMerger()

        for p in input_paths:
            try:
                merger.append(str(p))
            except Exception as e:
                return jsonify({'error': f'Could not read PDF: {p.name} — {str(e)}'}), 400

        merger.write(str(output_path))
        merger.close()

        if not output_path.exists() or output_path.stat().st_size == 0:
            return jsonify({'error': 'Merge produced an empty file'}), 500

        logger.info(f"Merge PDF successful: {len(input_paths)} files → {output_path.name}")
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name='merged.pdf',
            mimetype='application/pdf',
        )
    except Exception as e:
        logger.error(f"Merge PDF error: {e}", exc_info=True)
        return jsonify({'error': f'Merge failed: {str(e)}'}), 500
    finally:
        _cleanup(*input_paths)
        _cleanup(output_path)


# ─── Route: Compress PDF ─────────────────────────────────────────────────────

@app.route('/api/convert/compress-pdf', methods=['POST'])
def compress_pdf():
    input_path = output_path = None
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        if not _validate_extension(file.filename, ['.pdf']):
            return jsonify({'error': 'Invalid file type. Expected: .pdf'}), 400

        input_path = _save_upload(file, 'pdf')
        output_path = CONVERTED_DIR / _unique_name('.pdf')

        reader = PdfReader(str(input_path))
        writer = PdfWriter()

        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)

        # Remove metadata to reduce size
        writer.add_metadata({})

        with open(str(output_path), 'wb') as f_out:
            writer.write(f_out)

        original_size = input_path.stat().st_size
        compressed_size = output_path.stat().st_size
        ratio = ((original_size - compressed_size) / original_size * 100) if original_size > 0 else 0

        logger.info(
            f"Compress PDF successful: {original_size} → {compressed_size} bytes "
            f"({ratio:.1f}% reduction)"
        )

        original_stem = Path(secure_filename(file.filename)).stem
        return send_file(
            str(output_path),
            as_attachment=True,
            download_name=f"{original_stem}_compressed.pdf",
            mimetype='application/pdf',
        )
    except Exception as e:
        logger.error(f"Compress PDF error: {e}", exc_info=True)
        return jsonify({'error': f'Compression failed: {str(e)}'}), 500
    finally:
        _cleanup(input_path, output_path)


# ─── LibreOffice helper ──────────────────────────────────────────────────────

def _convert_with_libreoffice(input_path: Path, output_format: str) -> Path:
    """
    Convert a file using LibreOffice headless mode.
    Returns the path to the converted file.
    """
    if not LIBREOFFICE:
        raise RuntimeError('LibreOffice not found')

    out_dir = CONVERTED_DIR / uuid.uuid4().hex
    out_dir.mkdir(exist_ok=True)

    cmd = [
        LIBREOFFICE,
        '--headless',
        '--convert-to', output_format,
        '--outdir', str(out_dir),
        str(input_path),
    ]

    logger.info(f"Running LibreOffice: {' '.join(cmd)}")
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

    if result.returncode != 0:
        logger.error(f"LibreOffice error: {result.stderr}")
        shutil.rmtree(out_dir, ignore_errors=True)
        raise RuntimeError(f'LibreOffice conversion failed: {result.stderr}')

    # Find the output file
    output_files = list(out_dir.glob(f'*.{output_format}'))
    if not output_files:
        shutil.rmtree(out_dir, ignore_errors=True)
        raise RuntimeError('LibreOffice did not produce an output file')

    # Move to converted dir
    final_path = CONVERTED_DIR / _unique_name(f'.{output_format}')
    shutil.move(str(output_files[0]), str(final_path))
    shutil.rmtree(out_dir, ignore_errors=True)

    return final_path


# ─── Error Handlers ──────────────────────────────────────────────────────────

@app.errorhandler(413)
def file_too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 20MB.'}), 413


@app.errorhandler(404)
def not_found(e):
    return jsonify({'error': 'Endpoint not found'}), 404


@app.errorhandler(500)
def server_error(e):
    return jsonify({'error': 'Internal server error'}), 500


# ─── Run ──────────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    logger.info('🚀 FileForge server starting on http://localhost:5000')
    app.run(host='0.0.0.0', port=5000, debug=True)
